"""
IndAI — File Processor Service
Extracts text content from uploaded files for code scanning.

Supports:
- Code files (.py, .js, .java, .cpp, .json, .xml, .yaml, etc.)
- Text files (.txt, .md, .rtf)
- PDF (.pdf) (Text + OCR fallback)
- Images (.png, .jpg, .jpeg, .webp) (OCR)
- Word (.docx)
- PowerPoint (.pptx)
- Excel (.xlsx)
- CSV (.csv)
- OpenDocument (.odt, .ods, .odp)
"""

import csv
import io
import os
import zipfile
import logging

logger = logging.getLogger(__name__)


class FileProcessor:
    """Extracts text content from various file formats."""

    # Extensions that can be read directly as text
    TEXT_EXTENSIONS = {
        # Code files
        '.py', '.js', '.ts', '.jsx', '.tsx', '.java', '.cpp', '.c', '.h',
        '.hpp', '.cs', '.go', '.rb', '.rs', '.php', '.swift', '.kt', '.scala',
        '.r', '.m', '.sql', '.sh', '.bash', '.ps1', '.bat', '.cmd',
        '.html', '.htm', '.css', '.scss', '.less', '.sass',
        '.json', '.xml', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf',
        '.env', '.gitignore', '.dockerfile',
        # Text files
        '.txt', '.md', '.markdown', '.rst', '.rtf', '.log',
        # Data files
        '.csv',
    }

    # Binary formats requiring special processing
    BINARY_EXTENSIONS = {
        '.pdf', '.docx', '.pptx', '.xlsx',
        '.odt', '.ods', '.odp',
    }

    # Image formats for OCR
    IMAGE_EXTENSIONS = {
        '.png', '.jpg', '.jpeg', '.webp', '.bmp'
    }

    # Media files (not supported for text extraction)
    MEDIA_EXTENSIONS = {
        '.gif', '.svg',
        '.mp3', '.wav', '.m4a', '.ogg', '.flac',
        '.mp4', '.mov', '.avi', '.mkv', '.webm',
    }

    # Language detection from extension
    EXTENSION_TO_LANGUAGE = {
        '.py': 'python', '.pyw': 'python',
        '.js': 'javascript', '.jsx': 'javascript', '.mjs': 'javascript',
        '.ts': 'typescript', '.tsx': 'typescript',
        '.java': 'java',
        '.cpp': 'csharp', '.c': 'csharp', '.h': 'csharp', '.hpp': 'csharp',
        '.cs': 'csharp',
        '.go': 'go',
        '.rb': 'ruby',
        '.php': 'php',
        '.rs': 'rust',
        '.swift': 'swift',
    }

    # Max file size: 5MB for images/PDFs
    MAX_FILE_SIZE = 5 * 1024 * 1024

    @classmethod
    def process(cls, file_storage):
        """
        Process an uploaded file and extract text content.

        Args:
            file_storage: Flask FileStorage object

        Returns:
            dict: {
                'content': str,       # Extracted text
                'language': str,      # Detected language
                'filename': str,      # Original filename
                'file_type': str,     # Type category
                'lines': int,         # Line count
            }

        Raises:
            ValueError: If file is unsupported, too large, or empty.
        """
        filename = file_storage.filename or "unknown"
        ext = os.path.splitext(filename)[1].lower()

        # Check file size
        file_storage.seek(0, 2)  # Seek to end
        file_size = file_storage.tell()
        file_storage.seek(0)     # Reset to beginning

        if file_size > cls.MAX_FILE_SIZE:
            raise ValueError(
                f"File too large ({file_size // 1024}KB). "
                f"Maximum allowed size is {cls.MAX_FILE_SIZE // (1024*1024)}MB."
            )

        if file_size == 0:
            raise ValueError("File is empty.")

        # Route to appropriate handler
        if ext in cls.MEDIA_EXTENSIONS:
            raise ValueError(
                f"Media files ({ext}) cannot be scanned for code vulnerabilities. "
                f"Please paste the code content directly, or upload a code/document/image file."
            )

        if ext in cls.TEXT_EXTENSIONS:
            content = cls._read_text(file_storage)
            file_type = "code" if ext not in {'.txt', '.md', '.log', '.csv', '.rtf'} else "text"
        elif ext in cls.IMAGE_EXTENSIONS:
            content = cls._read_image_ocr(file_storage)
            file_type = "image_ocr"
        elif ext == '.pdf':
            content = cls._read_pdf(file_storage)
            file_type = "pdf"
        elif ext == '.docx':
            content = cls._read_docx(file_storage)
            file_type = "docx"
        elif ext == '.pptx':
            content = cls._read_pptx(file_storage)
            file_type = "pptx"
        elif ext == '.xlsx':
            content = cls._read_xlsx(file_storage)
            file_type = "xlsx"
        elif ext in {'.odt', '.ods', '.odp'}:
            content = cls._read_odf(file_storage, ext)
            file_type = "opendocument"
        else:
            # Try to read as text anyway
            try:
                content = cls._read_text(file_storage)
                file_type = "text"
            except Exception:
                raise ValueError(
                    f"Unsupported file type ({ext}). "
                    f"Please upload a code file, image, text file, or document."
                )

        content = content.strip()
        if not content:
            raise ValueError(
                f"No text content could be extracted from '{filename}'. "
                f"The file may be empty or the image quality is too low for OCR."
            )

        # Detect language from extension (default to JS if unknown)
        language = cls.EXTENSION_TO_LANGUAGE.get(ext, 'javascript')

        return {
            'content': content,
            'language': language,
            'filename': filename,
            'file_type': file_type,
            'lines': content.count('\n') + 1,
        }

    @classmethod
    def _read_text(cls, file_storage):
        """Read a text/code file."""
        raw = file_storage.read()
        # Try UTF-8 first, then latin-1 as fallback
        try:
            return raw.decode('utf-8')
        except UnicodeDecodeError:
            return raw.decode('latin-1')

    @classmethod
    def _read_image_ocr(cls, file_storage):
        """Extract text from an image using Tesseract OCR."""
        try:
            import pytesseract
            from PIL import Image
        except ImportError:
            raise ValueError("OCR dependencies (pytesseract, Pillow) are not installed on the server.")

        # Allow configuring Tesseract path via .env (e.g., if installed on Drive D)
        tesseract_cmd = os.environ.get("TESSERACT_CMD")
        if tesseract_cmd:
            # Strip surrounding quotes that dotenv may include
            tesseract_cmd = tesseract_cmd.strip('"').strip("'")
        if tesseract_cmd and os.path.exists(tesseract_cmd):
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        elif os.name == 'nt':
            # Check common Windows installation paths
            common_paths = [
                r'C:\Program Files\Tesseract-OCR\tesseract.exe',
                r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
                r'D:\Program Files\tesseract.exe',
                r'D:\Tesseract-OCR\tesseract.exe',
                r'D:\Program Files\Tesseract-OCR\tesseract.exe'
            ]
            for path in common_paths:
                if os.path.exists(path):
                    pytesseract.pytesseract.tesseract_cmd = path
                    break

        try:
            image = Image.open(file_storage)
            # You can configure tesseract options here if needed, e.g., --psm 6 for blocks of code
            text = pytesseract.image_to_string(image)
            return text
        except pytesseract.TesseractNotFoundError:
             logger.error("Tesseract executable not found.")
             raise ValueError("Tesseract-OCR is not installed or not found in PATH. Please install it or set TESSERACT_CMD in your .env file.")
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            raise ValueError("Failed to extract text from image. Make sure Tesseract-OCR is installed correctly.")

    @classmethod
    def _read_pdf(cls, file_storage):
        """Extract text from PDF using pdfplumber."""
        try:
            import pdfplumber
        except ImportError:
            # Fallback to PyPDF2 if pdfplumber isn't installed
            from PyPDF2 import PdfReader
            reader = PdfReader(file_storage)
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return '\n\n'.join(pages)

        pages = []
        try:
            with pdfplumber.open(file_storage) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        pages.append(text)
            return '\n\n'.join(pages)
        except Exception as e:
            logger.error(f"PDF extraction failed: {e}")
            raise ValueError("Failed to read PDF file. It might be corrupted or protected.")

    @classmethod
    def _read_docx(cls, file_storage):
        """Extract text from Word document."""
        from docx import Document

        doc = Document(file_storage)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n'.join(paragraphs)

    @classmethod
    def _read_pptx(cls, file_storage):
        """Extract text from PowerPoint presentation."""
        from pptx import Presentation

        prs = Presentation(file_storage)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                texts.append(f"--- Slide {slide_num} ---")
                texts.extend(slide_texts)
        return '\n'.join(texts)

    @classmethod
    def _read_xlsx(cls, file_storage):
        """Extract text from Excel spreadsheet."""
        from openpyxl import load_workbook

        wb = load_workbook(file_storage, read_only=True, data_only=True)
        texts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            texts.append(f"--- Sheet: {sheet} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                if row_text.strip():
                    texts.append(row_text)
        wb.close()
        return '\n'.join(texts)

    @classmethod
    def _read_odf(cls, file_storage, ext):
        """Extract text from OpenDocument Format files (.odt, .ods, .odp)."""
        # ODF files are ZIP archives containing content.xml
        raw = file_storage.read()
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as zf:
                if 'content.xml' in zf.namelist():
                    content_xml = zf.read('content.xml').decode('utf-8')
                    # Simple XML text extraction
                    import re
                    # Remove XML tags, keep text content
                    text = re.sub(r'<[^>]+>', ' ', content_xml)
                    # Clean up whitespace
                    text = re.sub(r'\s+', ' ', text).strip()
                    return text
                else:
                    raise ValueError("Invalid OpenDocument file.")
        except zipfile.BadZipFile:
            raise ValueError("File appears to be corrupted.")
